"""
Generate TruckFleet stakeholder presentation as a .pptx file.
Run: python3 scripts/generate_slides.py
Output: TruckFleet_Presentation.pptx (project root)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0F, 0x17, 0x2A)   # slide background
BLUE       = RGBColor(0x26, 0x5D, 0xD4)   # accent / headings
LIGHT_BLUE = RGBColor(0x60, 0x9C, 0xFF)   # sub-headings
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xA8, 0xB4, 0xC8)
GREEN      = RGBColor(0x34, 0xD3, 0x99)
ORANGE     = RGBColor(0xFB, 0x92, 0x3C)
RED        = RGBColor(0xF8, 0x71, 0x71)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=NAVY):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a plain rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def label(slide, text, left, top, width, height,
          font_size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    """Add a text box."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def accent_bar(slide, color=BLUE, height=0.07):
    """Thin horizontal accent bar at top of slide."""
    box(slide, 0, 0, 13.33, height, fill_color=color)


def section_heading(slide, title, subtitle=None):
    """Reusable heading block."""
    accent_bar(slide)
    label(slide, title,
          left=0.6, top=0.2, width=12, height=0.7,
          font_size=32, bold=True, color=WHITE)
    if subtitle:
        label(slide, subtitle,
              left=0.6, top=0.95, width=12, height=0.4,
              font_size=15, color=GRAY)


def bullet_block(slide, items, left, top, width=5.8, font_size=14,
                 color=WHITE, bullet="●", spacing=0.42):
    """Render a simple bullet list."""
    for i, item in enumerate(items):
        label(slide, f"{bullet}  {item}",
              left=left, top=top + i * spacing,
              width=width, height=0.4,
              font_size=font_size, color=color)


def card(slide, left, top, width, height, bg_color, title, body_lines,
         title_color=WHITE, body_color=GRAY, title_size=15, body_size=12):
    """Rounded-ish card (rectangle + text)."""
    box(slide, left, top, width, height, fill_color=bg_color)
    label(slide, title,
          left=left + 0.18, top=top + 0.12,
          width=width - 0.3, height=0.4,
          font_size=title_size, bold=True, color=title_color)
    for i, line in enumerate(body_lines):
        label(slide, line,
              left=left + 0.18, top=top + 0.52 + i * 0.33,
              width=width - 0.3, height=0.35,
              font_size=body_size, color=body_color)


# ════════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title / Cover"""
    s = blank_slide(prs)
    bg(s, NAVY)

    # large blue band on left
    box(s, 0, 0, 5.2, 7.5, fill_color=BLUE)

    # logo-style text
    label(s, "🚛", left=0.5, top=0.9, width=4, height=1.2,
          font_size=72, color=WHITE, align=PP_ALIGN.CENTER)
    label(s, "TruckFleet", left=0.2, top=2.1, width=4.8, height=0.9,
          font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    label(s, "Chemical Transport Management Platform",
          left=0.2, top=3.0, width=4.8, height=0.55,
          font_size=14, color=RGBColor(0xBD, 0xD7, 0xFF), align=PP_ALIGN.CENTER)

    # right side
    label(s, "Stakeholder Overview",
          left=5.6, top=2.6, width=7, height=0.7,
          font_size=30, bold=True, color=WHITE)
    label(s, "A web-based platform that unifies fleet operations,\n"
             "hazardous cargo compliance, and real-time driver\n"
             "coordination — in one system.",
          left=5.6, top=3.4, width=7.3, height=1.5,
          font_size=15, color=GRAY)
    label(s, "March 2026",
          left=5.6, top=6.5, width=4, height=0.4,
          font_size=12, color=GRAY, italic=True)


def slide_problem(prs):
    """Slide 2 — The Problem"""
    s = blank_slide(prs)
    bg(s, NAVY)
    section_heading(s, "The Problem", "Chemical transport operations are fragmented and error-prone")
    accent_bar(s)

    pain_points = [
        ("Compliance Risk",
         RED,
         ["No centralised cert tracking", "Manual hazard-class checks", "Missed HOS violations"]),
        ("Visibility Gaps",
         ORANGE,
         ["Dispatchers can't see trucks in real time", "No ETA alerts or offline warnings", "Phone calls for status updates"]),
        ("Paper-Heavy Ops",
         BLUE,
         ["Pre/post-trip forms on paper", "Proof of delivery lost or delayed", "SDS/MSDS documents not linked to trips"]),
    ]

    for i, (title, color, bullets) in enumerate(pain_points):
        left = 0.5 + i * 4.25
        box(s, left, 1.55, 3.9, 5.5, fill_color=RGBColor(0x15, 0x20, 0x38))
        # top colour strip
        box(s, left, 1.55, 3.9, 0.09, fill_color=color)
        label(s, title,
              left=left + 0.2, top=1.7, width=3.5, height=0.45,
              font_size=16, bold=True, color=color)
        bullet_block(s, bullets,
                     left=left + 0.15, top=2.3,
                     width=3.6, font_size=13, color=WHITE,
                     bullet="✗", spacing=0.55)


def slide_solution(prs):
    """Slide 3 — Our Solution"""
    s = blank_slide(prs)
    bg(s, NAVY)
    section_heading(s, "Our Solution", "TruckFleet — one platform for everyone in the operation")
    accent_bar(s)

    label(s, "TruckFleet is a role-aware web application that gives every stakeholder the exact tools "
             "they need — admins manage the fleet, dispatchers coordinate trips, and drivers execute deliveries "
             "— all with built-in compliance safeguards.",
          left=0.6, top=1.45, width=12.1, height=1.0,
          font_size=14, color=GRAY)

    roles = [
        ("👤  Admin", BLUE,
         ["Manage trucks & drivers", "Define chemical load library", "Enforce certification rules"]),
        ("📋  Dispatcher", GREEN,
         ["Create & assign trips", "Real-time fleet map", "Alert management"]),
        ("🚛  Driver", ORANGE,
         ["Mobile-first PWA", "Trip details & navigation", "Status updates & messaging"]),
    ]
    for i, (role, color, bullets) in enumerate(roles):
        left = 0.5 + i * 4.25
        box(s, left, 2.6, 3.9, 4.5, fill_color=RGBColor(0x15, 0x20, 0x38))
        label(s, role,
              left=left + 0.2, top=2.75, width=3.5, height=0.5,
              font_size=17, bold=True, color=color)
        bullet_block(s, bullets,
                     left=left + 0.2, top=3.35,
                     width=3.5, font_size=13, color=WHITE,
                     bullet="→", spacing=0.5)


def slide_features(prs):
    """Slide 4 — Key Features"""
    s = blank_slide(prs)
    bg(s, NAVY)
    section_heading(s, "Key Features", "Everything needed to run a compliant chemical fleet")
    accent_bar(s)

    features = [
        ("Fleet Management",    BLUE,   "Add, edit, and deactivate trucks with vehicle-type classification."),
        ("Driver Profiles",     BLUE,   "Track certifications (HazMat, TWIC, etc.) and license details."),
        ("Chemical Load Library", GREEN, "Define loads with UN numbers, hazard classes, and cert requirements."),
        ("Smart Dispatch",      GREEN,  "Auto-filters eligible drivers and trucks for each chemical type."),
        ("Real-Time Tracking",  ORANGE, "Live GPS map (Leaflet / OpenStreetMap) with offline alerts."),
        ("Mobile Driver PWA",   ORANGE, "Install-to-home-screen web app; no app store needed."),
        ("Digital Compliance",  RED,    "Pre/post-trip checklists, HOS logging, signature-based POD."),
        ("Messaging",           LIGHT_BLUE, "Per-trip dispatcher ↔ driver thread; no external tools needed."),
    ]

    cols = 4
    for i, (title, color, desc) in enumerate(features):
        col = i % cols
        row = i // cols
        left = 0.4 + col * 3.2
        top  = 1.55 + row * 2.6
        box(s, left, top, 3.0, 2.3, fill_color=RGBColor(0x15, 0x20, 0x38))
        box(s, left, top, 3.0, 0.07, fill_color=color)
        label(s, title,
              left=left + 0.15, top=top + 0.15, width=2.7, height=0.4,
              font_size=13, bold=True, color=color)
        label(s, desc,
              left=left + 0.15, top=top + 0.6, width=2.7, height=1.5,
              font_size=11, color=GRAY)


def slide_how_it_works(prs):
    """Slide 5 — How It Works (workflow)"""
    s = blank_slide(prs)
    bg(s, NAVY)
    section_heading(s, "How It Works", "A trip from creation to completion")
    accent_bar(s)

    steps = [
        ("1", "Define Load",      BLUE,   "Admin adds chemical with hazard class, UN number, cert requirements"),
        ("2", "Schedule Trip",    GREEN,  "Dispatcher picks load → system filters certified drivers & eligible trucks"),
        ("3", "Driver Notified",  ORANGE, "Driver receives assignment in mobile app with load details & navigation"),
        ("4", "Live Tracking",    LIGHT_BLUE, "GPS pings sent every 30 s → dispatcher sees live map with ETA"),
        ("5", "Proof of Delivery",WHITE, "Driver captures signature + photo → trip marked delivered"),
        ("6", "Compliance Log",   GRAY,  "Full audit trail: inspections, HOS, docs — visible to dispatcher"),
    ]

    arrow_top = 3.8
    step_w = 2.0
    gap = 0.22
    start_left = 0.35

    for i, (num, title, color, desc) in enumerate(steps):
        left = start_left + i * (step_w + gap)

        # circle
        circle = slide_add_circle(s, left + 0.6, 1.6, 0.8, 0.8, color)

        label(s, num,
              left=left + 0.6, top=1.6, width=0.8, height=0.8,
              font_size=22, bold=True, color=NAVY if color != NAVY else WHITE,
              align=PP_ALIGN.CENTER)

        label(s, title,
              left=left, top=2.55, width=step_w, height=0.45,
              font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)

        label(s, desc,
              left=left, top=3.05, width=step_w, height=1.4,
              font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

        # arrow between steps
        if i < len(steps) - 1:
            label(s, "→",
                  left=left + step_w + 0.02, top=1.75, width=0.25, height=0.5,
                  font_size=18, color=GRAY, align=PP_ALIGN.CENTER)


def slide_add_circle(slide, left_in, top_in, w_in, h_in, color):
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        9,  # OVAL
        Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def slide_architecture(prs):
    """Slide 6 — High-Level Architecture"""
    s = blank_slide(prs)
    bg(s, NAVY)
    section_heading(s, "High-Level Architecture", "Modern web stack — secure, scalable, zero-cost infrastructure")
    accent_bar(s)

    # three-layer diagram
    layers = [
        ("Frontend", BLUE,
         ["Next.js 16 (React 19)", "Tailwind CSS + shadcn/ui", "PWA (driver mobile)", "Leaflet map"]),
        ("Backend / API", GREEN,
         ["Next.js App Router", "BetterAuth (OAuth + Email)", "Drizzle ORM", "SSE real-time stream"]),
        ("Data & Storage", ORANGE,
         ["PostgreSQL", "Vercel Blob (files)", "Drizzle migrations", "Audit-ready schema"]),
    ]

    for i, (name, color, bullets) in enumerate(layers):
        left = 0.5 + i * 4.25
        box(s, left, 1.55, 3.9, 5.5, fill_color=RGBColor(0x10, 0x1B, 0x30))
        box(s, left, 1.55, 3.9, 0.09, fill_color=color)
        label(s, name,
              left=left + 0.2, top=1.7, width=3.5, height=0.45,
              font_size=16, bold=True, color=color)
        bullet_block(s, bullets,
                     left=left + 0.2, top=2.35,
                     width=3.5, font_size=13, color=WHITE,
                     bullet="▸", spacing=0.55)

    # arrows between layers
    for x in [4.42, 8.67]:
        label(s, "⟶", left=x, top=3.85, width=0.4, height=0.5,
              font_size=22, color=GRAY, align=PP_ALIGN.CENTER)


def slide_benefits(prs):
    """Slide 7 — Business Benefits"""
    s = blank_slide(prs)
    bg(s, NAVY)
    section_heading(s, "Business Benefits", "Quantifiable impact on operations and risk")
    accent_bar(s)

    benefits = [
        ("Reduced Compliance Risk", GREEN,
         "Automated cert checks prevent uncertified drivers from being assigned to hazardous loads."),
        ("Faster Dispatch Cycle", BLUE,
         "Smart filtering cuts trip-creation time; no manual cross-referencing of spreadsheets."),
        ("Real-Time Visibility", ORANGE,
         "Dispatchers see every truck's location and status without phone calls."),
        ("Paperless Operations", LIGHT_BLUE,
         "Digital checklists, e-signatures, and auto-linked docs eliminate paperwork delays."),
        ("Driver Productivity", GREEN,
         "Mobile PWA means drivers access everything on their phone — no laptop required."),
        ("Audit-Ready Records", ORANGE,
         "Every inspection, HOS log, and delivery signature is stored and searchable."),
    ]

    for i, (title, color, desc) in enumerate(benefits):
        col = i % 2
        row = i // 2
        left = 0.5 + col * 6.4
        top  = 1.55 + row * 1.9

        box(s, left, top, 6.1, 1.65, fill_color=RGBColor(0x15, 0x20, 0x38))
        box(s, left, top, 0.08, 1.65, fill_color=color)

        label(s, title,
              left=left + 0.25, top=top + 0.12, width=5.7, height=0.4,
              font_size=14, bold=True, color=color)
        label(s, desc,
              left=left + 0.25, top=top + 0.55, width=5.7, height=0.95,
              font_size=12, color=GRAY)


def slide_roadmap(prs):
    """Slide 8 — Development Roadmap"""
    s = blank_slide(prs)
    bg(s, NAVY)
    section_heading(s, "Development Roadmap", "Phased delivery — each phase delivers usable functionality")
    accent_bar(s)

    phases = [
        ("Phase 0", "Foundation",         "✅ Complete",  GREEN,
         "Auth, DB schema, role-based routing, seeded data"),
        ("Phase 1", "Fleet Management",   "✅ Complete",  GREEN,
         "Truck & driver CRUD, chemical load library, cert validation"),
        ("Phase 2", "Dispatch",           "✅ Complete",  GREEN,
         "Trip creation, lifecycle management, availability board"),
        ("Phase 3", "Driver Mobile",      "✅ Complete",  GREEN,
         "PWA, trip cards, MSDS viewer, dispatcher messaging"),
        ("Phase 4", "Live Tracking",      "✅ Complete",  GREEN,
         "GPS capture, Leaflet map, SSE stream, offline alerts"),
        ("Phase 5", "Compliance & POD",   "✅ Complete",  GREEN,
         "Inspections, HOS logging, e-signature, compliance reports"),
    ]

    for i, (phase, name, status, color, desc) in enumerate(phases):
        top = 1.55 + i * 0.93
        box(s, 0.5, top, 12.3, 0.82, fill_color=RGBColor(0x15, 0x20, 0x38))
        box(s, 0.5, top, 0.09, 0.82, fill_color=color)

        label(s, phase,
              left=0.7, top=top + 0.08, width=1.1, height=0.35,
              font_size=11, bold=True, color=color)
        label(s, name,
              left=1.85, top=top + 0.08, width=2.2, height=0.35,
              font_size=13, bold=True, color=WHITE)
        label(s, desc,
              left=4.1, top=top + 0.08, width=6.8, height=0.65,
              font_size=11, color=GRAY)
        label(s, status,
              left=10.95, top=top + 0.08, width=1.85, height=0.35,
              font_size=11, bold=True, color=color, align=PP_ALIGN.RIGHT)


def slide_next_steps(prs):
    """Slide 9 — Next Steps / Call to Action"""
    s = blank_slide(prs)
    bg(s, NAVY)
    accent_bar(s, color=GREEN)

    label(s, "Ready to Transform Your Fleet Operations?",
          left=1, top=1.4, width=11.3, height=1,
          font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    label(s, "TruckFleet is already running — Phase 0 complete, Phase 1 in active development.",
          left=1, top=2.55, width=11.3, height=0.55,
          font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

    next_steps = [
        ("Demo Access", BLUE,    "Request a live walkthrough of the current build"),
        ("Feedback",    GREEN,   "Prioritise upcoming phases based on your operational needs"),
        ("Timeline",    ORANGE,  "Agree a delivery schedule for Phases 1–5"),
    ]

    for i, (title, color, desc) in enumerate(next_steps):
        left = 1.1 + i * 3.8
        box(s, left, 3.4, 3.4, 2.5, fill_color=RGBColor(0x15, 0x20, 0x38))
        box(s, left, 3.4, 3.4, 0.09, fill_color=color)
        label(s, title,
              left=left + 0.2, top=3.58, width=3.0, height=0.45,
              font_size=16, bold=True, color=color)
        label(s, desc,
              left=left + 0.2, top=4.1, width=3.0, height=1.2,
              font_size=13, color=WHITE)

    label(s, "TruckFleet  |  2026",
          left=0, top=6.9, width=13.33, height=0.4,
          font_size=11, color=GRAY, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_features(prs)
    slide_how_it_works(prs)
    slide_architecture(prs)
    slide_benefits(prs)
    slide_roadmap(prs)
    slide_next_steps(prs)

    out = os.path.join(os.path.dirname(__file__), "..", "TruckFleet_Presentation.pptx")
    out = os.path.abspath(out)
    prs.save(out)
    print(f"✅  Saved: {out}")


if __name__ == "__main__":
    main()
